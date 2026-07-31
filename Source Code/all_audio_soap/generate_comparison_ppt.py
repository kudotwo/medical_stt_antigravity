"""
Generate Transcript Comparison PowerPoint
==========================================
Creates a PPT comparing 5 conversations:
  - Slide 1: Dialogue comparison (real txt vs STT output, first ~100 words)
  - Slide 2: Summary comparison (predicted vs ground truth SOAP summary, first ~100 words)
  - Slide 3: SOAP report comparison table (predicted vs ground truth)
"""

import csv
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
BASE_DIR = Path(__file__).parent
SOURCE_DIR = BASE_DIR.parent
AUDIO_DIR = SOURCE_DIR.parent / 'audio_sample' / 'audio_recordings' / 'Audio_Recordings'
TRANSCRIPT_DIR = SOURCE_DIR.parent / 'audio_sample' / 'audio_recordings' / 'Clean_Transcripts'
PREDICTED_DIR = BASE_DIR / 'predicted'
GROUND_TRUTH_DIR = BASE_DIR / 'ground_truth'

OUTPUT_FILE = BASE_DIR / 'Transcript_Comparison_Presentation.pptx'

# 5 cases from different specialties
CASES = [
    {'id': 'CAR0001', 'label': 'Cardiology — Chest Pain'},
    {'id': 'MSK0005', 'label': 'Musculoskeletal — Elbow Pain'},
    {'id': 'RES0050', 'label': 'Respiratory — Chronic Cough'},
    {'id': 'GAS0003', 'label': 'Gastroenterology — Abdominal Pain'},
    {'id': 'DER0001', 'label': 'Dermatology — Skin Rash'},
]

# ---------------------------------------------------------------------------
# Color palette — professional dark theme
# ---------------------------------------------------------------------------
BG_DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # Dark navy background
BG_CARD      = RGBColor(0x22, 0x22, 0x3A)   # Slightly lighter card bg
ACCENT_BLUE  = RGBColor(0x4E, 0xC9, 0xF0)   # Bright cyan accent
ACCENT_GREEN = RGBColor(0x50, 0xE3, 0xA0)   # Green for ground truth
ACCENT_ORANGE= RGBColor(0xF0, 0x9E, 0x4E)   # Orange for predicted
TEXT_WHITE   = RGBColor(0xF0, 0xF0, 0xF5)   # Near-white text
TEXT_GRAY    = RGBColor(0xA0, 0xA0, 0xB0)   # Muted gray text
TABLE_HEADER = RGBColor(0x2A, 0x2A, 0x4A)   # Table header bg
TABLE_ROW1   = RGBColor(0x1E, 0x1E, 0x36)   # Table row alternating 1
TABLE_ROW2   = RGBColor(0x26, 0x26, 0x40)   # Table row alternating 2
BORDER_COLOR = RGBColor(0x3A, 0x3A, 0x5A)   # Subtle border

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def read_soap_csv(filepath):
    """Read a SOAP CSV and return dict."""
    with open(filepath, 'r', encoding='utf-8') as f:
        reader = csv.DictReader(f)
        rows = list(reader)
        return rows[0] if rows else {}


def read_diarized_transcript(case_id):
    """Read the Gemini-diarized transcript and convert to D:/P: format."""
    txt_path = AUDIO_DIR / f'{case_id}_transcript.txt'
    if not txt_path.exists():
        return ''
    raw = txt_path.read_text(encoding='utf-8')
    lines = []
    for line in raw.strip().split('\n'):
        line = line.strip()
        if not line:
            continue
        # Strip timestamp prefix like "[0.0s->12.0s] "
        if '] ' in line:
            line = line.split('] ', 1)[1]
        # Convert "Doctor: ..." -> "D: ..." and "Patient: ..." -> "P: ..."
        if line.startswith('Doctor: '):
            line = 'D: ' + line[8:]
        elif line.startswith('Patient: '):
            line = 'P: ' + line[9:]
        lines.append(line)
    return '\n'.join(lines)


def read_clean_transcript(case_id):
    """Read clean transcript txt file (already has D:/P: format)."""
    txt_path = TRANSCRIPT_DIR / f'{case_id}.txt'
    if not txt_path.exists():
        return ''
    raw = txt_path.read_text(encoding='utf-8')
    # Clean up — keep only non-empty lines
    lines = [l.strip() for l in raw.split('\n') if l.strip()]
    return '\n'.join(lines)


def first_n_words(text, n=100):
    """Return the first n words of text, keeping line structure."""
    words = text.split()
    if len(words) <= n:
        return text
    # Take first n words, then find the natural line break
    truncated_text = ' '.join(words[:n])
    # Try to end at a line boundary for cleaner display
    lines = text.split('\n')
    result_lines = []
    word_count = 0
    for line in lines:
        line_words = len(line.split())
        if word_count + line_words > n:
            break
        result_lines.append(line)
        word_count += line_words
    if result_lines:
        return '\n'.join(result_lines) + '\n...'
    return truncated_text + '...'


def set_slide_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=16,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text_box(slide, left, top, width, height, lines, font_size=14,
                           color=TEXT_WHITE, font_name='Calibri'):
    """Add a text box with multiple lines (list of tuples: (text, bold, color))."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, txt_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = txt_color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox


# SOAP fields for the table
SOAP_FIELDS = [
    ('Chief Complaint', 'subjective.chief_complaint'),
    ('Symptoms', 'subjective.symptoms'),
    ('Symptom Onset', 'subjective.symptom_onset'),
    ('Medical History', 'subjective.medical_history'),
    ('Medications', 'subjective.current_medications'),
    ('Allergies', 'subjective.allergies'),
    ('Exam Findings', 'objective.physical_exam_findings'),
    ('Diagnosis', 'assessment.diagnosis'),
    ('Diff. Diagnosis', 'assessment.differential_diagnosis'),
    ('Prescribed Meds', 'plan.prescribed_medications'),
    ('Treatment Plan', 'plan.treatment_plan'),
    ('Investigations', 'plan.investigations_ordered'),
    ('Follow Up', 'plan.follow_up'),
]

def truncate_cell(text, max_chars=120):
    """Truncate text for table cells."""
    if not text or text.strip() == '':
        return '—'
    if len(text) > max_chars:
        return text[:max_chars] + '...'
    return text


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title_slide(prs):
    """Build the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, BG_DARK)

    # Title
    add_text_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
                 'Medical STT Pipeline', font_size=44, color=ACCENT_BLUE,
                 bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')

    # Subtitle
    add_text_box(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.8),
                 'Transcript Comparison: Predicted (Speech-to-Text) vs Ground Truth',
                 font_size=24, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    # Pipeline description
    add_text_box(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.6),
                 'Pipeline: Audio (.mp3) → Whisper Large-v3 (faster-whisper, GPU) → Gemini 2.0 Flash → SOAP Report',
                 font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Case count
    add_text_box(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.5),
                 '5 Representative Conversations from 5 Medical Specialties',
                 font_size=18, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # Bottom info
    add_text_box(slide, Inches(2), Inches(6.2), Inches(9.3), Inches(0.4),
                 'Dataset: MTS-Dialog (Medical Transcription Scoring) | 271 audio files processed',
                 font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


def build_case_header_slide(prs, case_num, case_label, case_id):
    """Build a section divider slide for each case."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Case number badge
    add_rounded_rect(slide, Inches(5.2), Inches(2.0), Inches(2.9), Inches(0.7), ACCENT_BLUE)
    add_text_box(slide, Inches(5.2), Inches(2.05), Inches(2.9), Inches(0.6),
                 f'Conversation {case_num} of 5', font_size=20, color=BG_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Case label
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1.0),
                 case_label, font_size=36, color=TEXT_WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Case ID
    add_text_box(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
                 f'Case ID: {case_id}', font_size=20, color=TEXT_GRAY,
                 alignment=PP_ALIGN.CENTER)


def build_dialogue_slide(prs, case_id, case_label, stt_text, gt_text):
    """Slide 1: Dialogue comparison — first 100 words each."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title bar
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
                 f'Dialogue Comparison — {case_label}',
                 font_size=24, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.35),
                 'First ~100 words of each transcript shown below',
                 font_size=14, color=TEXT_GRAY)

    # Left panel: Ground Truth
    left_x = Inches(0.4)
    panel_w = Inches(6.0)
    panel_top = Inches(1.5)
    panel_h = Inches(5.5)

    add_rounded_rect(slide, left_x, panel_top, panel_w, panel_h, BG_CARD, BORDER_COLOR)

    # GT label
    add_rounded_rect(slide, left_x + Inches(0.3), panel_top + Inches(0.2),
                     Inches(3.0), Inches(0.45), ACCENT_GREEN)
    add_text_box(slide, left_x + Inches(0.3), panel_top + Inches(0.22),
                 Inches(3.0), Inches(0.4),
                 '📄 Ground Truth Transcript', font_size=14, color=BG_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # GT text
    gt_100 = first_n_words(gt_text, 100)
    add_text_box(slide, left_x + Inches(0.3), panel_top + Inches(0.85),
                 panel_w - Inches(0.6), panel_h - Inches(1.1),
                 gt_100, font_size=14, color=TEXT_WHITE)

    # Right panel: STT (Predicted)
    right_x = Inches(6.9)

    add_rounded_rect(slide, right_x, panel_top, panel_w, panel_h, BG_CARD, BORDER_COLOR)

    # STT label
    add_rounded_rect(slide, right_x + Inches(0.3), panel_top + Inches(0.2),
                     Inches(3.4), Inches(0.45), ACCENT_ORANGE)
    add_text_box(slide, right_x + Inches(0.3), panel_top + Inches(0.22),
                 Inches(3.4), Inches(0.4),
                 '🎙️ Speech-to-Text Output', font_size=14, color=BG_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # STT text
    stt_100 = first_n_words(stt_text, 100)
    add_text_box(slide, right_x + Inches(0.3), panel_top + Inches(0.85),
                 panel_w - Inches(0.6), panel_h - Inches(1.1),
                 stt_100, font_size=14, color=TEXT_WHITE)


def build_summary_slide(prs, case_id, case_label, pred_summary, gt_summary):
    """Slide 2: Summary comparison — first 100 words each."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
                 f'Summary Comparison — {case_label}',
                 font_size=24, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.35),
                 'SOAP summary generated from each transcript (first ~100 words)',
                 font_size=14, color=TEXT_GRAY)

    # Left panel: Ground Truth Summary
    left_x = Inches(0.4)
    panel_w = Inches(6.0)
    panel_top = Inches(1.5)
    panel_h = Inches(5.5)

    add_rounded_rect(slide, left_x, panel_top, panel_w, panel_h, BG_CARD, BORDER_COLOR)

    add_rounded_rect(slide, left_x + Inches(0.3), panel_top + Inches(0.2),
                     Inches(4.0), Inches(0.45), ACCENT_GREEN)
    add_text_box(slide, left_x + Inches(0.3), panel_top + Inches(0.22),
                 Inches(4.0), Inches(0.4),
                 '📄 Summary from Ground Truth', font_size=14, color=BG_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)

    gt_100 = first_n_words(gt_summary, 100)
    add_text_box(slide, left_x + Inches(0.3), panel_top + Inches(0.85),
                 panel_w - Inches(0.6), panel_h - Inches(1.1),
                 gt_100, font_size=14, color=TEXT_WHITE)

    # Right panel: Predicted Summary
    right_x = Inches(6.9)

    add_rounded_rect(slide, right_x, panel_top, panel_w, panel_h, BG_CARD, BORDER_COLOR)

    add_rounded_rect(slide, right_x + Inches(0.3), panel_top + Inches(0.2),
                     Inches(4.0), Inches(0.45), ACCENT_ORANGE)
    add_text_box(slide, right_x + Inches(0.3), panel_top + Inches(0.22),
                 Inches(4.0), Inches(0.4),
                 '🎙️ Summary from STT Output', font_size=14, color=BG_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)

    pred_100 = first_n_words(pred_summary, 100)
    add_text_box(slide, right_x + Inches(0.3), panel_top + Inches(0.85),
                 panel_w - Inches(0.6), panel_h - Inches(1.1),
                 pred_100, font_size=14, color=TEXT_WHITE)


def build_soap_table_slide(prs, case_id, case_label, pred_data, gt_data):
    """Slide 3: SOAP report table — predicted vs ground truth."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Title
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5),
                 f'SOAP Report Comparison — {case_label}',
                 font_size=22, color=ACCENT_BLUE, bold=True)

    # Build table
    rows = len(SOAP_FIELDS) + 1  # +1 for header
    cols = 3  # Field | Ground Truth | Predicted
    table_left = Inches(0.3)
    table_top = Inches(0.85)
    table_width = Inches(12.7)
    table_height = Inches(6.4)

    table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(5.45)
    table.columns[2].width = Inches(5.45)

    # Style header row
    headers = ['SOAP Field', 'Ground Truth', 'Predicted (STT)']
    for j, header_text in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.name = 'Calibri'
            if j == 1:
                paragraph.font.color.rgb = ACCENT_GREEN
            elif j == 2:
                paragraph.font.color.rgb = ACCENT_ORANGE
            else:
                paragraph.font.color.rgb = ACCENT_BLUE
            paragraph.alignment = PP_ALIGN.CENTER

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Populate data rows
    for i, (field_label, field_key) in enumerate(SOAP_FIELDS, start=1):
        gt_val = gt_data.get(field_key, '')
        pred_val = pred_data.get(field_key, '')

        row_bg = TABLE_ROW1 if i % 2 == 1 else TABLE_ROW2

        # Field name cell
        cell0 = table.cell(i, 0)
        cell0.text = field_label
        cell0.fill.solid()
        cell0.fill.fore_color.rgb = row_bg
        for p in cell0.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
            p.font.name = 'Calibri'
        cell0.vertical_anchor = MSO_ANCHOR.TOP

        # Ground truth cell
        cell1 = table.cell(i, 1)
        cell1.text = truncate_cell(gt_val, 140)
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = row_bg
        for p in cell1.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_WHITE
            p.font.name = 'Calibri'
        cell1.vertical_anchor = MSO_ANCHOR.TOP

        # Predicted cell
        cell2 = table.cell(i, 2)
        cell2.text = truncate_cell(pred_val, 140)
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = row_bg
        for p in cell2.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_WHITE
            p.font.name = 'Calibri'
        cell2.vertical_anchor = MSO_ANCHOR.TOP

    # Remove table borders for cleaner look (set to thin subtle lines)
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    build_title_slide(prs)

    # Process each case
    for case_num, case in enumerate(CASES, 1):
        case_id = case['id']
        case_label = case['label']

        print(f'Processing {case_id} ({case_label})...')

        # Read data
        stt_text = read_diarized_transcript(case_id)
        gt_text = read_clean_transcript(case_id)
        pred_data = read_soap_csv(PREDICTED_DIR / f'{case_id}_soap.csv')
        gt_data = read_soap_csv(GROUND_TRUTH_DIR / f'{case_id}_soap.csv')

        pred_summary = pred_data.get('summary', '')
        gt_summary = gt_data.get('summary', '')

        # Build section divider
        build_case_header_slide(prs, case_num, case_label, case_id)

        # Slide 1: Dialogue comparison
        build_dialogue_slide(prs, case_id, case_label, stt_text, gt_text)

        # Slide 2: Summary comparison
        build_summary_slide(prs, case_id, case_label, pred_summary, gt_summary)

        # Slide 3: SOAP table
        build_soap_table_slide(prs, case_id, case_label, pred_data, gt_data)

    # Save
    prs.save(str(OUTPUT_FILE))
    print(f'\n✅ Presentation saved to: {OUTPUT_FILE}')
    print(f'   Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()

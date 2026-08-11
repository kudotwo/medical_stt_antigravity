"""
Generate RES0050 Gemini 3.6 Flash Evaluation Deck
=================================================
Creates a 3-main-slide presentation for conversation RES0050 (Respiratory — Chronic Cough / COPD):
  1. Slide 1: Speech-to-Text transcript comparison with real transcript (first ~100 words)
  2. Slide 2: Transcript summary comparison (first ~100 words)
  3. Slide 3: SOAP report comparison table (Ground Truth vs Gemini 3.6 Flash Predicted)
"""

import csv
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
BASE_DIR = Path(__file__).parent
SOURCE_DIR = BASE_DIR.parent
AUDIO_DIR = SOURCE_DIR.parent / 'audio_sample' / 'audio_recordings' / 'Audio_Recordings'
TRANSCRIPT_DIR = SOURCE_DIR.parent / 'audio_sample' / 'audio_recordings' / 'Clean_Transcripts'
PREDICTED_DIR = BASE_DIR / 'predicted'
GROUND_TRUTH_DIR = BASE_DIR / 'ground_truth'

OUTPUT_FILE = BASE_DIR / 'RES0050_Gemini_3.6_Flash_Presentation.pptx'

CASE_ID = 'RES0050'
CASE_LABEL = 'RES0050: Respiratory — Chronic Cough / COPD (Gemini 3.6 Flash)'

# ---------------------------------------------------------------------------
# Color palette — dark navy theme
# ---------------------------------------------------------------------------
BG_DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # Dark navy background
BG_CARD      = RGBColor(0x22, 0x22, 0x3A)   # Card background
ACCENT_BLUE  = RGBColor(0x4E, 0xC9, 0xF0)   # Cyan accent
ACCENT_GREEN = RGBColor(0x50, 0xE3, 0xA0)   # Green for ground truth
ACCENT_ORANGE= RGBColor(0xF0, 0x9E, 0x4E)   # Orange for predicted
TEXT_WHITE   = RGBColor(0xF0, 0xF0, 0xF5)   # White text
TEXT_GRAY    = RGBColor(0xA0, 0xA0, 0xB0)   # Gray text
TABLE_HEADER = RGBColor(0x2A, 0x2A, 0x4A)   # Table header bg
TABLE_ROW1   = RGBColor(0x1E, 0x1E, 0x36)   # Alternating row 1
TABLE_ROW2   = RGBColor(0x26, 0x26, 0x40)   # Alternating row 2
BORDER_COLOR = RGBColor(0x3A, 0x3A, 0x5A)   # Subtle border

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def read_soap_csv(filepath):
    """Read a SOAP CSV and return dict."""
    if not filepath.exists():
        return {}
    with open(filepath, 'r', encoding='utf-8') as f:
        reader = csv.DictReader(f)
        rows = list(reader)
        return rows[0] if rows else {}


def read_diarized_transcript(case_id):
    """Read the Gemini-diarized transcript and convert to D:/P: format."""
    txt_path = AUDIO_DIR / f'{case_id}_transcript.txt'
    if not txt_path.exists():
        return ''
    raw = txt_path.read_text(encoding='utf-8')
    lines = []
    for line in raw.strip().split('\n'):
        line = line.strip()
        if not line:
            continue
        if '] ' in line:
            line = line.split('] ', 1)[1]
        if line.startswith('Doctor: '):
            line = 'D: ' + line[8:]
        elif line.startswith('Patient: '):
            line = 'P: ' + line[9:]
        lines.append(line)
    return '\n'.join(lines)


def read_clean_transcript(case_id):
    """Read clean transcript txt file (already has D:/P: format)."""
    txt_path = TRANSCRIPT_DIR / f'{case_id}.txt'
    if not txt_path.exists():
        return ''
    raw = txt_path.read_text(encoding='utf-8')
    lines = [l.strip() for l in raw.split('\n') if l.strip()]
    return '\n'.join(lines)


def first_n_words(text, n=100):
    """Return the first n words of text, keeping line structure."""
    words = text.split()
    if len(words) <= n:
        return text
    truncated_text = ' '.join(words[:n])
    lines = text.split('\n')
    result_lines = []
    word_count = 0
    for line in lines:
        line_words = len(line.split())
        if word_count + line_words > n:
            break
        result_lines.append(line)
        word_count += line_words
    if result_lines:
        return '\n'.join(result_lines) + '\n...'
    return truncated_text + '...'


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=16,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


SOAP_FIELDS = [
    ('Chief Complaint', 'subjective.chief_complaint'),
    ('Symptoms', 'subjective.symptoms'),
    ('Symptom Onset', 'subjective.symptom_onset'),
    ('Medical History', 'subjective.medical_history'),
    ('Medications', 'subjective.current_medications'),
    ('Allergies', 'subjective.allergies'),
    ('Exam Findings', 'objective.physical_exam_findings'),
    ('Diagnosis', 'assessment.diagnosis'),
    ('Diff. Diagnosis', 'assessment.differential_diagnosis'),
    ('Prescribed Meds', 'plan.prescribed_medications'),
    ('Treatment Plan', 'plan.treatment_plan'),
    ('Investigations', 'plan.investigations_ordered'),
    ('Follow Up', 'plan.follow_up'),
]


def truncate_cell(text, max_chars=120):
    if not text or text.strip() == '':
        return '—'
    if len(text) > max_chars:
        return text[:max_chars] + '...'
    return text

# ---------------------------------------------------------------------------
# Slide Builders
# ---------------------------------------------------------------------------

def build_dialogue_slide(prs, stt_text, gt_text):
    """Slide 1: Speech-to-Text Transcript Comparison with Real Transcript (first 100 words)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Header
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
                 f'Slide 1: Speech-to-Text vs Real Transcript — {CASE_LABEL}',
                 font_size=24, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.35),
                 'Comparison of Ground Truth dialogue vs faster-whisper + Gemini 3.6 Flash STT output (first ~100 words)',
                 font_size=14, color=TEXT_GRAY)

    # Left panel: Real Ground Truth Transcript
    left_x = Inches(0.4)
    panel_w = Inches(6.0)
    panel_top = Inches(1.5)
    panel_h = Inches(5.5)

    add_rounded_rect(slide, left_x, panel_top, panel_w, panel_h, BG_CARD, BORDER_COLOR)
    add_rounded_rect(slide, left_x + Inches(0.3), panel_top + Inches(0.2),
                     Inches(3.5), Inches(0.45), ACCENT_GREEN)
    add_text_box(slide, left_x + Inches(0.3), panel_top + Inches(0.22),
                 Inches(3.5), Inches(0.4),
                 '📄 Real Transcript (Ground Truth)', font_size=14, color=BG_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)

    gt_100 = first_n_words(gt_text, 100)
    add_text_box(slide, left_x + Inches(0.3), panel_top + Inches(0.85),
                 panel_w - Inches(0.6), panel_h - Inches(1.1),
                 gt_100, font_size=14, color=TEXT_WHITE)

    # Right panel: Speech-to-Text Output
    right_x = Inches(6.9)

    add_rounded_rect(slide, right_x, panel_top, panel_w, panel_h, BG_CARD, BORDER_COLOR)
    add_rounded_rect(slide, right_x + Inches(0.3), panel_top + Inches(0.2),
                     Inches(3.8), Inches(0.45), ACCENT_ORANGE)
    add_text_box(slide, right_x + Inches(0.3), panel_top + Inches(0.22),
                 Inches(3.8), Inches(0.4),
                 '🎙️ Speech-to-Text Transcript (Whisper + Gemini 3.6)', font_size=14, color=BG_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)

    stt_100 = first_n_words(stt_text, 100)
    add_text_box(slide, right_x + Inches(0.3), panel_top + Inches(0.85),
                 panel_w - Inches(0.6), panel_h - Inches(1.1),
                 stt_100, font_size=14, color=TEXT_WHITE)


def build_summary_slide(prs, pred_summary, gt_summary):
    """Slide 2: Transcript Summary Comparison (first 100 words)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Header
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
                 f'Slide 2: Transcript Summary Comparison — {CASE_LABEL}',
                 font_size=24, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.35),
                 'Clinical summary extracted from Ground Truth vs Gemini 3.6 Flash STT output (first ~100 words)',
                 font_size=14, color=TEXT_GRAY)

    # Left panel: Ground Truth Summary
    left_x = Inches(0.4)
    panel_w = Inches(6.0)
    panel_top = Inches(1.5)
    panel_h = Inches(5.5)

    add_rounded_rect(slide, left_x, panel_top, panel_w, panel_h, BG_CARD, BORDER_COLOR)
    add_rounded_rect(slide, left_x + Inches(0.3), panel_top + Inches(0.2),
                     Inches(4.0), Inches(0.45), ACCENT_GREEN)
    add_text_box(slide, left_x + Inches(0.3), panel_top + Inches(0.22),
                 Inches(4.0), Inches(0.4),
                 '📄 Ground Truth Summary', font_size=14, color=BG_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)

    gt_100 = first_n_words(gt_summary, 100)
    add_text_box(slide, left_x + Inches(0.3), panel_top + Inches(0.85),
                 panel_w - Inches(0.6), panel_h - Inches(1.1),
                 gt_100, font_size=14, color=TEXT_WHITE)

    # Right panel: Predicted Summary
    right_x = Inches(6.9)

    add_rounded_rect(slide, right_x, panel_top, panel_w, panel_h, BG_CARD, BORDER_COLOR)
    add_rounded_rect(slide, right_x + Inches(0.3), panel_top + Inches(0.2),
                     Inches(4.2), Inches(0.45), ACCENT_ORANGE)
    add_text_box(slide, right_x + Inches(0.3), panel_top + Inches(0.22),
                 Inches(4.2), Inches(0.4),
                 '🎙️ Gemini 3.6 Flash Summary (STT)', font_size=14, color=BG_DARK,
                 bold=True, alignment=PP_ALIGN.CENTER)

    pred_100 = first_n_words(pred_summary, 100)
    add_text_box(slide, right_x + Inches(0.3), panel_top + Inches(0.85),
                 panel_w - Inches(0.6), panel_h - Inches(1.1),
                 pred_100, font_size=14, color=TEXT_WHITE)


def build_soap_table_slide(prs, pred_data, gt_data):
    """Slide 3: SOAP Report Comparison Table (Ground Truth vs Gemini 3.6 Flash Predicted)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Header
    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5),
                 f'Slide 3: Structured SOAP Report Comparison — {CASE_LABEL}',
                 font_size=22, color=ACCENT_BLUE, bold=True)

    rows = len(SOAP_FIELDS) + 1
    cols = 3
    table_left = Inches(0.3)
    table_top = Inches(0.85)
    table_width = Inches(12.7)
    table_height = Inches(6.4)

    table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(5.45)
    table.columns[2].width = Inches(5.45)

    headers = ['SOAP Field', 'Ground Truth', 'Predicted (Gemini 3.6 Flash)']
    for j, header_text in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.name = 'Calibri'
            if j == 1:
                p.font.color.rgb = ACCENT_GREEN
            elif j == 2:
                p.font.color.rgb = ACCENT_ORANGE
            else:
                p.font.color.rgb = ACCENT_BLUE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, (field_label, field_key) in enumerate(SOAP_FIELDS, start=1):
        gt_val = gt_data.get(field_key, '')
        pred_val = pred_data.get(field_key, '')

        row_bg = TABLE_ROW1 if i % 2 == 1 else TABLE_ROW2

        # Field Label
        cell0 = table.cell(i, 0)
        cell0.text = field_label
        cell0.fill.solid()
        cell0.fill.fore_color.rgb = row_bg
        for p in cell0.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
            p.font.name = 'Calibri'
        cell0.vertical_anchor = MSO_ANCHOR.TOP

        # Ground Truth
        cell1 = table.cell(i, 1)
        cell1.text = truncate_cell(gt_val, 140)
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = row_bg
        for p in cell1.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_WHITE
            p.font.name = 'Calibri'
        cell1.vertical_anchor = MSO_ANCHOR.TOP

        # Predicted
        cell2 = table.cell(i, 2)
        cell2.text = truncate_cell(pred_val, 140)
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = row_bg
        for p in cell2.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT_WHITE
            p.font.name = 'Calibri'
        cell2.vertical_anchor = MSO_ANCHOR.TOP

    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    stt_text = read_diarized_transcript(CASE_ID)
    gt_text = read_clean_transcript(CASE_ID)
    pred_data = read_soap_csv(PREDICTED_DIR / f'{CASE_ID}_soap.csv')
    gt_data = read_soap_csv(GROUND_TRUTH_DIR / f'{CASE_ID}_soap.csv')

    pred_summary = pred_data.get('summary', '')
    gt_summary = gt_data.get('summary', '')

    print(f'Building RES0050 Gemini 3.6 Flash Presentation...')

    # 3 Main Slides
    build_dialogue_slide(prs, stt_text, gt_text)
    build_summary_slide(prs, pred_summary, gt_summary)
    build_soap_table_slide(prs, pred_data, gt_data)

    prs.save(str(OUTPUT_FILE))
    print(f'✅ Presentation saved successfully to: {OUTPUT_FILE}')
    print(f'   Total main slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()

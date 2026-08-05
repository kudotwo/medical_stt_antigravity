"""
generate_model_comparison_deck.py
=================================
Creates a presentation comparing all three SOAP extraction models:
  - gemini-3.1-flash-lite
  - gemini-3.5-flash
  - gpt-5.6-terra

Slides:
  1.  Title / Overview
  2.  Overall average performance (all 3 models, 3 metrics, 0-100%)
  3.  Full results heatmap table (all 15 data points)
  4-8. Per-case clustered bar charts (CAR0001 → DER0001)

Metrics (paper-aligned, 0-100% scale):
  - Clinical Fidelity
  - Logical Consistency
  - SOAP Structural Fidelity  (renamed from predictive_accuracy)

Judge: gemini-3.5-flash-lite
"""

import csv
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── Paths ────────────────────────────────────────────────────────────────────
BASE_DIR    = Path(__file__).parent
METRICS_CSV = BASE_DIR / 'model_comparison_metrics.csv'
OUTPUT_FILE = BASE_DIR / 'Model_Comparison_3Models_Presentation.pptx'

CASES = [
    {'id': 'CAR0001', 'label': 'Cardiology — Chest Pain'},
    {'id': 'MSK0005', 'label': 'Musculoskeletal — Elbow Pain'},
    {'id': 'RES0050', 'label': 'Respiratory — Chronic Cough'},
    {'id': 'GAS0003', 'label': 'Gastroenterology — Abdominal Pain'},
    {'id': 'DER0001', 'label': 'Dermatology — Skin Rash'},
]

MODEL_LABELS = {
    'gemini-3.1-flash-lite': 'Gemini 3.1 Flash Lite',
    'gemini-3.5-flash':      'Gemini 3.5 Flash',
    'gpt-5.6-terra':         'GPT-5.6 Terra',
}

# ── Color palette ─────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0F, 0x0F, 0x23)
BG_CARD       = RGBColor(0x1A, 0x1A, 0x35)
ACCENT_CYAN   = RGBColor(0x00, 0xD4, 0xFF)   # Gemini 3.1 Lite
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)   # Gemini 3.5
ACCENT_PURPLE = RGBColor(0xAA, 0x5C, 0xF7)   # GPT-5.6 Terra
ACCENT_GREEN  = RGBColor(0x00, 0xE5, 0x96)
TEXT_WHITE    = RGBColor(0xF0, 0xF0, 0xFF)
TEXT_GRAY     = RGBColor(0x88, 0x88, 0xAA)
BORDER_DIM    = RGBColor(0x2A, 0x2A, 0x48)
HEADER_BG     = RGBColor(0x20, 0x20, 0x40)
ROW1          = RGBColor(0x18, 0x18, 0x32)
ROW2          = RGBColor(0x20, 0x20, 0x3C)

MODEL_COLORS = [ACCENT_CYAN, ACCENT_ORANGE, ACCENT_PURPLE]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
METRICS = ['clinical_fidelity', 'logical_consistency', 'soap_structural_fidelity']
METRIC_LABELS = ['Clinical Fidelity', 'Logical Consistency', 'SOAP Structural Fidelity']


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, l, t, w, h, text, sz=16, color=TEXT_WHITE, bold=False,
       align=PP_ALIGN.LEFT, name='Calibri'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return box

def add_chart(slide, chart_data, l, t, w, h, title, y_max=100):
    gf    = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, chart_data)
    chart = gf.chart

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
    chart.chart_title.text_frame.paragraphs[0].font.bold = True

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.color.rgb = TEXT_WHITE

    va = chart.value_axis
    va.maximum_scale = float(y_max)
    va.minimum_scale = 0.0
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = RGBColor(0x30, 0x30, 0x50)
    va.tick_labels.font.color.rgb = TEXT_GRAY
    va.tick_labels.font.size = Pt(10)

    ca = chart.category_axis
    ca.tick_labels.font.color.rgb = TEXT_WHITE
    ca.tick_labels.font.size = Pt(11)

    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = MODEL_COLORS[i % len(MODEL_COLORS)]

    return chart


# ── Slide 1: Title ────────────────────────────────────────────────────────────

def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Decorative top bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_CYAN
    bar.line.fill.background()

    tb(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
       'Medical STT — SOAP Extraction', sz=38, color=ACCENT_CYAN, bold=True,
       align=PP_ALIGN.CENTER)
    tb(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.7),
       '3-Model Performance Comparison', sz=26, color=TEXT_WHITE, bold=True,
       align=PP_ALIGN.CENTER)
    tb(slide, Inches(1.5), Inches(3.55), Inches(10.3), Inches(0.5),
       'Gemini 3.1 Flash Lite  ·  Gemini 3.5 Flash  ·  GPT-5.6 Terra',
       sz=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    tb(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.4),
       '5 Clinical Conversations  |  3 Metrics  |  Judge: gemini-3.5-flash-lite',
       sz=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    tb(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.35),
       'Metrics: Clinical Fidelity · Logical Consistency · SOAP Structural Fidelity  (0–100%)',
       sz=13, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    # Bottom bar
    bot = slide.shapes.add_shape(1, Inches(0), Inches(7.3), SLIDE_W, Inches(0.12))
    bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT_PURPLE
    bot.line.fill.background()


# ── Slide 2: Overall average bar chart ────────────────────────────────────────

def build_overall_slide(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6),
       'Overall Model Performance', sz=30, color=ACCENT_GREEN, bold=True)
    tb(slide, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.35),
       'Average scores across all 5 clinical conversations  (0–100%)',
       sz=14, color=TEXT_GRAY)

    avg_df = df.groupby('model')[METRICS].mean().reset_index()

    chart_data = CategoryChartData()
    chart_data.categories = METRIC_LABELS
    for _, row in avg_df.iterrows():
        label = MODEL_LABELS.get(row['model'], row['model'])
        chart_data.add_series(label, tuple(round(row[m], 1) for m in METRICS))

    add_chart(slide, chart_data, Inches(1.0), Inches(1.55), Inches(11.3), Inches(5.5),
              'Average Metric Scores (0–100%)')


# ── Slide 3: Full results heatmap table ───────────────────────────────────────

def build_heatmap_slide(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55),
       'Full Results — All 15 Evaluations', sz=26, color=ACCENT_CYAN, bold=True)
    tb(slide, Inches(0.4), Inches(0.82), Inches(12.5), Inches(0.35),
       '5 cases × 3 models  |  Metric scores 0–100%  |  Judge: gemini-3.5-flash-lite',
       sz=13, color=TEXT_GRAY)

    rows = len(df) + 1  # +1 header
    cols = 6            # Case | Model | CF | LC | SSF | Avg
    t_left = Inches(0.3)
    t_top  = Inches(1.3)
    t_w    = Inches(12.7)
    t_h    = Inches(5.9)

    tbl_shape = slide.shapes.add_table(rows, cols, t_left, t_top, t_w, t_h)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.9)
    tbl.columns[1].width = Inches(3.2)
    tbl.columns[2].width = Inches(1.9)
    tbl.columns[3].width = Inches(2.0)
    tbl.columns[4].width = Inches(2.3)
    tbl.columns[5].width = Inches(1.4)

    headers = ['Case', 'Model', 'Clinical\nFidelity', 'Logical\nConsistency',
               'SOAP Structural\nFidelity', 'Avg']
    hdr_colors = [TEXT_WHITE, TEXT_WHITE, ACCENT_CYAN, ACCENT_ORANGE, ACCENT_PURPLE, ACCENT_GREEN]

    for j, (h, hc) in enumerate(zip(headers, hdr_colors)):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10); p.font.bold = True
            p.font.color.rgb = hc; p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Sort for nice grouping by case
    df_sorted = df.sort_values(['case_id', 'model']).reset_index(drop=True)

    for i, (_, row) in enumerate(df_sorted.iterrows(), start=1):
        bg = ROW1 if i % 2 == 0 else ROW2
        cf  = int(row['clinical_fidelity'])
        lc  = int(row['logical_consistency'])
        ssf = int(row['soap_structural_fidelity'])
        avg_val = round((cf + lc + ssf) / 3, 1)

        values = [
            row['case_id'],
            MODEL_LABELS.get(row['model'], row['model']),
            f"{cf}%",
            f"{lc}%",
            f"{ssf}%",
            f"{avg_val}%",
        ]
        val_colors = [TEXT_WHITE, TEXT_WHITE, ACCENT_CYAN, ACCENT_ORANGE, ACCENT_PURPLE, ACCENT_GREEN]

        for j, (val, vc) in enumerate(zip(values, val_colors)):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10); p.font.name = 'Calibri'
                p.font.color.rgb = vc
                p.font.bold = (j == 5)
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for j in range(cols):
            cell = tbl.cell(i, j)
            cell.margin_left  = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top   = Inches(0.03)
            cell.margin_bottom = Inches(0.03)


# ── Slides 4-8: Per-case bar charts ──────────────────────────────────────────

def build_case_slide(prs, df, case):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.55),
       f"Case: {case['label']}", sz=26, color=ACCENT_CYAN, bold=True)
    tb(slide, Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.35),
       f"Conversation {case['id']}  |  Performance metrics 0–100%",
       sz=14, color=TEXT_GRAY)

    case_df = df[df['case_id'] == case['id']]
    chart_data = CategoryChartData()
    chart_data.categories = METRIC_LABELS
    for _, row in case_df.iterrows():
        label = MODEL_LABELS.get(row['model'], row['model'])
        chart_data.add_series(label, tuple(int(row[m]) for m in METRICS))

    add_chart(slide, chart_data, Inches(1.0), Inches(1.55), Inches(11.3), Inches(5.5),
              f"Metric Scores — {case['id']}  (0–100%)")


# ── Slide 9: Model averages summary table ─────────────────────────────────────

def build_summary_table_slide(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6),
       'Model Performance Summary', sz=30, color=ACCENT_GREEN, bold=True)
    tb(slide, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.35),
       'Average scores across all 5 cases  (0–100%)  |  Paper-aligned metrics',
       sz=14, color=TEXT_GRAY)

    avg_df = df.groupby('model')[METRICS].mean().round(1).reset_index()
    avg_df['overall'] = avg_df[METRICS].mean(axis=1).round(1)
    avg_df = avg_df.sort_values('overall', ascending=False).reset_index(drop=True)

    rows = len(avg_df) + 1
    cols = 5
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(2.0), Inches(1.7), Inches(9.3), Inches(2.8))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.8)
    tbl.columns[1].width = Inches(1.9)
    tbl.columns[2].width = Inches(1.9)
    tbl.columns[3].width = Inches(2.1)
    tbl.columns[4].width = Inches(1.6)

    headers    = ['Model', 'Clinical\nFidelity', 'Logical\nConsistency', 'SOAP Structural\nFidelity', 'Overall\nAvg']
    hdr_colors = [TEXT_WHITE, ACCENT_CYAN, ACCENT_ORANGE, ACCENT_PURPLE, ACCENT_GREEN]
    for j, (h, hc) in enumerate(zip(headers, hdr_colors)):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12); p.font.bold = True
            p.font.color.rgb = hc; p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, (_, row) in enumerate(avg_df.iterrows(), start=1):
        bg = ROW1 if i % 2 == 1 else ROW2
        vals = [
            MODEL_LABELS.get(row['model'], row['model']),
            f"{row['clinical_fidelity']}%",
            f"{row['logical_consistency']}%",
            f"{row['soap_structural_fidelity']}%",
            f"{row['overall']}%",
        ]
        val_colors = [TEXT_WHITE, ACCENT_CYAN, ACCENT_ORANGE, ACCENT_PURPLE, ACCENT_GREEN]
        for j, (v, vc) in enumerate(zip(vals, val_colors)):
            cell = tbl.cell(i, j)
            cell.text = str(v)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13); p.font.name = 'Calibri'
                p.font.color.rgb = vc
                p.font.bold = (j in (0, 4))
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left  = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top   = Inches(0.06)
            cell.margin_bottom = Inches(0.06)

    # Key finding callout
    winner = avg_df.iloc[0]
    winner_label = MODEL_LABELS.get(winner['model'], winner['model'])
    tb(slide, Inches(2.0), Inches(4.8), Inches(9.3), Inches(0.55),
       f'🏆  Highest overall average: {winner_label}  ({winner["overall"]}%)',
       sz=16, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

    tb(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.8),
       'Note: All three models perform exceptionally well (>94% average), demonstrating '
       'that modern LLMs are highly capable at structured SOAP extraction from STT transcripts.',
       sz=13, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    if not METRICS_CSV.exists():
        print(f'Error: {METRICS_CSV} not found. Run evaluate_models.py + rejudge_metrics.py first.')
        return

    df = pd.read_csv(METRICS_CSV)
    if 'soap_structural_fidelity' not in df.columns:
        print('Error: CSV uses old metric names. Run rejudge_metrics.py first to get paper-aligned metrics.')
        return

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print('Building 3-Model Comparison Presentation...')
    print('  Slide 1: Title')
    build_title_slide(prs)

    print('  Slide 2: Overall average performance')
    build_overall_slide(prs, df)

    print('  Slide 3: Full results heatmap table')
    build_heatmap_slide(prs, df)

    for case in CASES:
        print(f"  Slide: {case['id']} per-case bar chart")
        build_case_slide(prs, df, case)

    print('  Final slide: Summary table')
    build_summary_table_slide(prs, df)

    prs.save(str(OUTPUT_FILE))
    print(f'\n✅ Presentation saved → {OUTPUT_FILE}')
    print(f'   Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()

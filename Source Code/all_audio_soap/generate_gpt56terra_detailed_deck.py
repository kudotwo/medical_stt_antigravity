"""
generate_gpt56terra_detailed_deck.py
=====================================
Creates a 4-slide-per-case detailed presentation for GPT-5.6 Terra SOAP extraction.
Covers all 5 clinical conversations (20 slides total).

Per-case slides:
  1. Title / case overview with scores
  2. Speech-to-Text transcript: Ground Truth vs GPT-5.6 Terra diarized
  3. Clinical summary: Ground Truth vs GPT-5.6 Terra
  4. Structured SOAP report: Ground Truth vs GPT-5.6 Terra (table)
"""

import csv
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE_DIR         = Path(__file__).parent
SOURCE_DIR       = BASE_DIR.parent
TRANSCRIPT_DIR   = SOURCE_DIR.parent / 'audio_sample' / 'audio_recordings' / 'Clean_Transcripts'
PREDICTED_DIR    = BASE_DIR / 'predicted' / 'gpt-5.6-terra'
GROUND_TRUTH_DIR = BASE_DIR / 'ground_truth'
METRICS_DIR      = BASE_DIR / 'metrics_output'
OUTPUT_FILE      = BASE_DIR / 'GPT56Terra_Detailed_Presentation.pptx'

CASES = ['CAR0001', 'MSK0005', 'RES0050', 'GAS0003', 'DER0001']
CASE_LABELS = {
    'CAR0001': 'Cardiology — Chest Pain',
    'MSK0005': 'Musculoskeletal — Elbow Pain',
    'RES0050': 'Respiratory — Chronic Cough',
    'GAS0003': 'Gastroenterology — Abdominal Pain',
    'DER0001': 'Dermatology — Skin Rash',
}

MODEL_ID    = 'gpt-5.6-terra'
MODEL_LABEL = 'GPT-5.6 Terra'
SAFE_MODEL  = 'gpt-5.6-terra'

# ── Color palette — GPT/OpenAI theme ─────────────────────────────────────────
BG_DARK      = RGBColor(0x0A, 0x0A, 0x1A)
BG_CARD      = RGBColor(0x14, 0x14, 0x2E)
ACCENT_PURP  = RGBColor(0xAA, 0x5C, 0xF7)   # GPT brand purple
ACCENT_GREEN = RGBColor(0x00, 0xE5, 0x96)   # Ground truth
ACCENT_GOLD  = RGBColor(0xFF, 0xCC, 0x00)   # Score highlight
TEXT_WHITE   = RGBColor(0xF0, 0xF0, 0xFF)
TEXT_GRAY    = RGBColor(0x88, 0x88, 0xAA)
BORDER_DIM   = RGBColor(0x28, 0x28, 0x48)
HEADER_BG    = RGBColor(0x1E, 0x1E, 0x40)
ROW1         = RGBColor(0x16, 0x16, 0x30)
ROW2         = RGBColor(0x1E, 0x1E, 0x3A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SOAP_FIELDS = [
    ('Chief Complaint',  'subjective.chief_complaint'),
    ('Symptoms',         'subjective.symptoms'),
    ('Symptom Onset',    'subjective.symptom_onset'),
    ('Medical History',  'subjective.medical_history'),
    ('Curr. Medications','subjective.current_medications'),
    ('Allergies',        'subjective.allergies'),
    ('Exam Findings',    'objective.physical_exam_findings'),
    ('Diagnosis',        'assessment.diagnosis'),
    ('Diff. Diagnosis',  'assessment.differential_diagnosis'),
    ('Prescribed Meds',  'plan.prescribed_medications'),
    ('Treatment Plan',   'plan.treatment_plan'),
    ('Investigations',   'plan.investigations_ordered'),
    ('Follow Up',        'plan.follow_up'),
]


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def tb(slide, l, t, w, h, text, sz=14, color=TEXT_WHITE, bold=False,
       align=PP_ALIGN.LEFT, name='Calibri'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return box

def card(slide, l, t, w, h, fill=BG_CARD, border=BORDER_DIM):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.adjustments[0] = 0.04
    return s

def pill(slide, l, t, w, h, fill, text, text_color=BG_DARK, sz=13):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.adjustments[0] = 0.5
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = text_color; p.font.bold = True
    p.font.name = 'Calibri'; p.alignment = PP_ALIGN.CENTER
    return s

def read_soap_csv(path):
    if not path.exists(): return {}
    with open(path, 'r', encoding='utf-8') as f:
        rows = list(csv.DictReader(f))
        return rows[0] if rows else {}

def read_diarized(case_id):
    path = PREDICTED_DIR / f'{case_id}_transcript.txt'
    if not path.exists(): return ''
    raw = path.read_text(encoding='utf-8')
    lines = []
    for line in raw.strip().split('\n'):
        line = line.strip()
        if not line: continue
        if '] ' in line: line = line.split('] ', 1)[1]
        if line.startswith('Doctor: '):   line = 'D: ' + line[8:]
        elif line.startswith('Patient: '): line = 'P: ' + line[9:]
        lines.append(line)
    return '\n'.join(lines)

def read_gt_transcript(case_id):
    path = TRANSCRIPT_DIR / f'{case_id}.txt'
    if not path.exists(): return ''
    return '\n'.join(l.strip() for l in path.read_text(encoding='utf-8').split('\n') if l.strip())

def first_n_words(text, n=100):
    if not text: return ''
    words = text.split()
    if len(words) <= n: return text
    lines, count = [], 0
    for line in text.split('\n'):
        lw = len(line.split())
        if count + lw > n: break
        lines.append(line); count += lw
    return ('\n'.join(lines) + '\n...') if lines else ' '.join(words[:n]) + '...'

def truncate_cell(text, max_chars=130):
    if not text or str(text).strip() in ('', 'nan'): return '—'
    text = str(text)
    return text[:max_chars] + '...' if len(text) > max_chars else text

def load_scores(case_id):
    path = METRICS_DIR / f'{case_id}_{SAFE_MODEL}.json'
    if not path.exists(): return {}
    with open(path) as f: return json.load(f)


# ── Slide builders ────────────────────────────────────────────────────────────

def build_case_title_slide(prs, case_id):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_PURP
    bar.line.fill.background()

    tb(slide, Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.75),
       f'Case Study: {case_id}', sz=42, color=ACCENT_PURP, bold=True,
       align=PP_ALIGN.CENTER)
    tb(slide, Inches(1.0), Inches(1.85), Inches(11.3), Inches(0.5),
       CASE_LABELS.get(case_id, case_id), sz=22, color=TEXT_WHITE,
       align=PP_ALIGN.CENTER)
    tb(slide, Inches(1.0), Inches(2.45), Inches(11.3), Inches(0.4),
       f'SOAP Extraction Analysis  ·  {MODEL_LABEL}', sz=16, color=TEXT_GRAY,
       align=PP_ALIGN.CENTER)

    # Score pills
    scores = load_scores(case_id)
    if scores:
        cf  = scores.get('clinical_fidelity', '—')
        lc  = scores.get('logical_consistency', '—')
        ssf = scores.get('soap_structural_fidelity', '—')
        avg = round((cf + lc + ssf) / 3, 1) if all(isinstance(x, (int, float)) for x in [cf, lc, ssf]) else '—'

        pill_y = Inches(3.3)
        pill_h = Inches(0.85)
        labels = [
            ('Clinical Fidelity',       f'{cf}%',  ACCENT_PURP),
            ('Logical Consistency',     f'{lc}%',  ACCENT_PURP),
            ('SOAP Structural Fidelity',f'{ssf}%', ACCENT_PURP),
            ('Overall Average',         f'{avg}%', ACCENT_GOLD),
        ]
        total_w = Inches(11.0)
        pill_w  = Inches(2.5)
        gap     = (total_w - 4 * pill_w) / 3
        start_x = Inches(1.2)

        for i, (label, val, color) in enumerate(labels):
            x = start_x + i * (pill_w + gap)
            card(slide, x, pill_y, pill_w, pill_h, BG_CARD, BORDER_DIM)
            tb(slide, x, pill_y + Inches(0.08), pill_w, Inches(0.35),
               label, sz=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
            tb(slide, x, pill_y + Inches(0.42), pill_w, Inches(0.42),
               val, sz=24, color=color, bold=True, align=PP_ALIGN.CENTER)

    tb(slide, Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.35),
       f'Evaluated by gemini-3.5-flash-lite  |  Scale: 0–100%  |  Paper-aligned metrics',
       sz=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Bottom bar
    bot = slide.shapes.add_shape(1, Inches(0), Inches(7.3), SLIDE_W, Inches(0.1))
    bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT_PURP
    bot.line.fill.background()


def build_transcript_slide(prs, case_id):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
       f'{case_id} — Speech-to-Text Transcript Comparison',
       sz=22, color=ACCENT_PURP, bold=True)
    tb(slide, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.32),
       f'Ground Truth dialogue vs {MODEL_LABEL} diarized STT output (first ~100 words)',
       sz=13, color=TEXT_GRAY)

    panel_w = Inches(6.1)
    panel_t = Inches(1.28)
    panel_h = Inches(5.85)

    # Left — Ground Truth
    card(slide, Inches(0.3), panel_t, panel_w, panel_h)
    pill(slide, Inches(0.55), panel_t + Inches(0.18), Inches(3.8), Inches(0.42),
         ACCENT_GREEN, '📄  Real Transcript (Ground Truth)')
    gt_text = first_n_words(read_gt_transcript(case_id), 100)
    tb(slide, Inches(0.55), panel_t + Inches(0.75), panel_w - Inches(0.45), panel_h - Inches(1.0),
       gt_text, sz=13, color=TEXT_WHITE)

    # Right — GPT
    rx = Inches(6.75)
    card(slide, rx, panel_t, panel_w, panel_h)
    pill(slide, rx + Inches(0.25), panel_t + Inches(0.18), Inches(4.1), Inches(0.42),
         ACCENT_PURP, f'🤖  {MODEL_LABEL} Diarized STT', text_color=TEXT_WHITE)
    gpt_text = first_n_words(read_diarized(case_id), 100)
    tb(slide, rx + Inches(0.25), panel_t + Inches(0.75), panel_w - Inches(0.45), panel_h - Inches(1.0),
       gpt_text, sz=13, color=TEXT_WHITE)


def build_summary_slide(prs, case_id):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
       f'{case_id} — Clinical Summary Comparison',
       sz=22, color=ACCENT_PURP, bold=True)
    tb(slide, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.32),
       f'AI-extracted clinical summary: Ground Truth vs {MODEL_LABEL} (first ~100 words)',
       sz=13, color=TEXT_GRAY)

    panel_w = Inches(6.1)
    panel_t = Inches(1.28)
    panel_h = Inches(5.85)

    gt_data   = read_soap_csv(GROUND_TRUTH_DIR / f'{case_id}_soap.csv')
    pred_data = read_soap_csv(PREDICTED_DIR    / f'{case_id}_soap.csv')

    card(slide, Inches(0.3), panel_t, panel_w, panel_h)
    pill(slide, Inches(0.55), panel_t + Inches(0.18), Inches(3.6), Inches(0.42),
         ACCENT_GREEN, '📄  Ground Truth Summary')
    tb(slide, Inches(0.55), panel_t + Inches(0.75), panel_w - Inches(0.45), panel_h - Inches(1.0),
       first_n_words(gt_data.get('summary', 'N/A'), 100), sz=13, color=TEXT_WHITE)

    rx = Inches(6.75)
    card(slide, rx, panel_t, panel_w, panel_h)
    pill(slide, rx + Inches(0.25), panel_t + Inches(0.18), Inches(3.8), Inches(0.42),
         ACCENT_PURP, f'🤖  {MODEL_LABEL} Summary', text_color=TEXT_WHITE)
    tb(slide, rx + Inches(0.25), panel_t + Inches(0.75), panel_w - Inches(0.45), panel_h - Inches(1.0),
       first_n_words(pred_data.get('summary', 'N/A'), 100), sz=13, color=TEXT_WHITE)


def build_soap_table_slide(prs, case_id):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    tb(slide, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.5),
       f'{case_id} — Structured SOAP Report Comparison',
       sz=20, color=ACCENT_PURP, bold=True)

    gt_data   = read_soap_csv(GROUND_TRUTH_DIR / f'{case_id}_soap.csv')
    pred_data = read_soap_csv(PREDICTED_DIR    / f'{case_id}_soap.csv')

    rows = len(SOAP_FIELDS) + 1
    cols = 3
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(0.3), Inches(0.78), Inches(12.7), Inches(6.55))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.85)
    tbl.columns[1].width = Inches(5.43)
    tbl.columns[2].width = Inches(5.42)

    headers = ['SOAP Field', 'Ground Truth', f'Predicted ({MODEL_LABEL})']
    hdr_clrs = [ACCENT_PURP, ACCENT_GREEN, ACCENT_PURP]
    for j, (h, hc) in enumerate(zip(headers, hdr_clrs)):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11); p.font.bold = True
            p.font.color.rgb = hc; p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, (label, key) in enumerate(SOAP_FIELDS, start=1):
        bg = ROW1 if i % 2 == 1 else ROW2
        gt_v   = truncate_cell(gt_data.get(key, ''))
        pred_v = truncate_cell(pred_data.get(key, ''))

        for j, (val, vc) in enumerate([(label, ACCENT_PURP), (gt_v, TEXT_WHITE), (pred_v, TEXT_WHITE)]):
            cell = tbl.cell(i, j)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = vc
                p.font.bold = (j == 0)
                p.font.name = 'Calibri'
            cell.vertical_anchor = MSO_ANCHOR.TOP
            cell.margin_left  = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top   = Inches(0.03)
            cell.margin_bottom = Inches(0.03)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print(f'Building {MODEL_LABEL} Detailed Presentation...')

    for case_id in CASES:
        print(f'  Adding 4 slides for {case_id}...')
        build_case_title_slide(prs, case_id)
        build_transcript_slide(prs, case_id)
        build_summary_slide(prs, case_id)
        build_soap_table_slide(prs, case_id)

    prs.save(str(OUTPUT_FILE))
    print(f'\n✅ Presentation saved → {OUTPUT_FILE}')
    print(f'   Total slides: {len(prs.slides)}  ({len(CASES)} cases × 4 slides)')


if __name__ == '__main__':
    main()
